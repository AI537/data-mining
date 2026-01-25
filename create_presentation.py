
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_pptx():
    prs = Presentation()

    # Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Used Car Sales Data Mining Project"
    subtitle.text = "Analysis of Price, Mileage, and Market Segments"

    # 2. Agenda
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Agenda"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Project Overview"
    p = tf.add_paragraph()
    p.text = "Dataset & Preparation"
    p = tf.add_paragraph()
    p.text = "Methodology"
    p = tf.add_paragraph()
    p.text = "Clustering Analysis"
    p = tf.add_paragraph()
    p.text = "Anomaly Detection"
    p = tf.add_paragraph()
    p.text = "Association Rules"
    p = tf.add_paragraph()
    p.text = "Conclusion"

    # 3. Introduction
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Introduction & Domain"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Domain: Automotive E-commerce"
    p = tf.add_paragraph()
    p.text = "Objective: Apply data mining to uncover pricing patterns and market segments."
    p = tf.add_paragraph()
    p.text = "Key Questions:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "What drives car prices?"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Can we group similar cars?"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Are there fraud/error anomalies?"

    # 4. Dataset Description
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Dataset Description"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Source: 'used_car_sales.csv'"
    p = tf.add_paragraph()
    p.text = "Size: ~10,000 records (Sampled for analysis)"
    p = tf.add_paragraph()
    p.text = "Key Features:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Numerical: Price, Mileage, Year, Engine HP"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Categorical: Fuel, Gearbox, Car Type"

    # 5. Data Preparation
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Data Preparation Process"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Cleaning:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Renamed columns to standard format."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Filtered invalid data (Price <= 0, Mileage <= 0)."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Removed duplicates and missing values."
    p = tf.add_paragraph()
    p.text = "Transformation:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Scaled numerical features (StandardScaler)."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "One-Hot Encoded categorical variables."

    # 6. Clustering Results
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Clustering (K-Means)"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Goal: Market Segmentation"
    p = tf.add_paragraph()
    p.text = "Algorithm: K-Means (k=5)"
    p = tf.add_paragraph()
    p.text = "Findings:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Identified distinct groups (e.g., Budget-High Mileage vs Luxury-New)."
    
    # Add Image if exists
    img_path = 'cluster_plot.png'
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(3.5)
        height = Inches(3.5)
        slide.shapes.add_picture(img_path, left, top, height=height)

    # 7. Anomaly Detection
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Outlier Detection"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Goal: Detect Anomalies/Fraud"
    p = tf.add_paragraph()
    p.text = "Algorithm: Isolation Forest"
    p = tf.add_paragraph()
    p.text = "Findings:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Detected ~5% outliers."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Mostly high-price cars with high mileage (potential errors)."
    
    # Add Image if exists
    img_path = 'anomaly_plot.png'
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(3.5)
        height = Inches(3.5)
        slide.shapes.add_picture(img_path, left, top, height=height)

    # 8. Association Rules
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Association Rules"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Goal: Find Relationships"
    p = tf.add_paragraph()
    p.text = "Algorithm: Apriori"
    p = tf.add_paragraph()
    p.text = "Key Implications:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Features like 'Automatic' + 'SUV' strongly imply 'High Price'."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Helps in understanding feature bundles."

    # Add Image if exists
    img_path = 'association_plot.png'
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(3.5)
        height = Inches(3.5)
        slide.shapes.add_picture(img_path, left, top, height=height)

    # 9. Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Summary:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Successfully processed and analyzed car sales data."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Generated a runnable Pipeline for future data."
    p = tf.add_paragraph()
    p.text = "Impact:"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Data-driven pricing strategies."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "Better inventory management."

    # Save
    prs.save('DataMining_Presentation.pptx')
    print("Presentation generated successfully.")

if __name__ == "__main__":
    create_pptx()
